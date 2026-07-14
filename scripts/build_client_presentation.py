from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "SwimCRM_client_presentation_ru.pptx"
ASSETS = ROOT / "docs" / "presentation_assets"

NAVY = RGBColor(25, 38, 52)
BODY = RGBColor(73, 84, 97)
MUTED = RGBColor(111, 126, 140)
TEAL = RGBColor(0, 129, 146)
TEAL_DARK = RGBColor(0, 91, 106)
TEAL_SOFT = RGBColor(226, 247, 250)
BLUE_SOFT = RGBColor(232, 241, 255)
GREEN_SOFT = RGBColor(231, 248, 239)
AMBER_SOFT = RGBColor(255, 245, 219)
RED_SOFT = RGBColor(255, 236, 232)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(216, 226, 232)


SLIDES = [
    {
        "title": "SwimCRM для школы плавания",
        "subtitle": "Операционная CRM, которая связывает расписание, абонементы, оплаты и коммуникацию с родителями.",
        "tag": "Клиентская презентация | июль 2026",
        "image": "cover.png",
    },
    {
        "title": "О проекте",
        "body": [
            "SwimCRM — веб-система для частной школы плавания.",
            "Она помогает администрации вести клиентов, группы, тренеров, расписание, посещаемость, абонементы и оплаты в одном месте.",
            "Основные пользователи: администратор, тренер и родитель/клиент.",
        ],
        "image": "about.png",
    },
    {
        "title": "Ключевые возможности сейчас",
        "cards": [
            ("Основной функционал", "Клиенты и семьи, группы, тренеры, расписание, посещаемость, абонементы, начисления и платежи."),
            ("Автоматизация", "Списание занятий по статусам, контроль конфликтов тренера, отчеты, импорт/экспорт, планировщик уведомлений."),
            ("Интерфейс", "React SPA для админа, клиентский и тренерский контуры, адаптивная верстка и единая дизайн-система."),
        ],
        "image": "capabilities.png",
    },
    {
        "title": "Как это работает",
        "body": [
            "Администратор создает клиента или семейный аккаунт.",
            "Назначает участника в группу и выпускает абонемент.",
            "Система создает начисление и фиксирует оплату или чек.",
            "Тренер отмечает посещаемость после занятия.",
            "CRM обновляет остаток занятий, задолженности и напоминания.",
        ],
        "image": "flow.png",
    },
    {
        "title": "Примеры использования",
        "cards": [
            ("Продажа абонемента", "Админ создает абонемент, начисление и платеж. Клиент сразу видит актуальный статус."),
            ("Контроль посещаемости", "Тренер отмечает занятие, CRM сама списывает или не списывает занятие по бизнес-правилам."),
            ("Работа с долгами", "Админ видит просрочки, клиентов с малым остатком занятий и может запускать напоминания."),
            ("Отчетность", "Руководитель получает доходы, посещаемость и неоплаченные счета без ручной сборки таблиц."),
        ],
        "image": "cases.png",
    },
    {
        "title": "Ограничения — честно",
        "body": [
            "Онлайн-оплата пока не подключена: платежи вводятся администратором или подтверждаются по загруженному чеку.",
            "Telegram и SMS работают через provider adapters; для боевой доставки нужны реальные провайдерные ключи и контрольные отправки.",
            "Планировщик уведомлений готов как management command и Celery task; в production его нужно включить в Celery beat или cron.",
            "PostgreSQL-релизный контур покрыт тестами, backup/restore проверкой и production env-check.",
        ],
        "image": "limits.png",
    },
    {
        "title": "Что можно улучшить",
        "cards": [
            ("UX", "Упростить массовые действия в админке и добавить быстрые фильтры по группам, оплатам и остаткам занятий."),
            ("Функционал", "Настроить провайдерные аккаунты Telegram/SMS, расширить массовые рассылки и админ-экшены для импортов и отчетов."),
            ("Автоматизация", "Включить Celery/Redis в production, регулярные restore drills и мониторинг ошибок уведомлений."),
            ("Коммерция", "Добавить онлайн-оплаты и автоматическое сопоставление платежей, когда бизнес будет готов."),
        ],
        "image": "improvements.png",
    },
    {
        "title": "Roadmap развития",
        "cards": [
            ("Quick wins", "Production scheduler, мониторинг фоновых задач, админские действия для отчетов."),
            ("Среднесрочно", "Провайдерные аккаунты Telegram/SMS, массовые рассылки, расширенные UX-сценарии для администратора и тренера."),
            ("Продвинутые функции", "Онлайн-оплата, сверка банковских платежей, BI-дашборды и прогноз продлений."),
        ],
        "image": "roadmap.png",
    },
    {
        "title": "Бизнес-ценность",
        "body": [
            "Меньше ручной сверки оплат, посещаемости и остатков занятий.",
            "Быстрее видно долги, окончания абонементов и клиентов для продления.",
            "История действий, RODO-согласия, отчеты и прозрачные правила списания повышают управляемость бизнеса.",
        ],
        "image": "value.png",
    },
    {
        "title": "Итог",
        "body": [
            "SwimCRM уже является практичной основой для операционного управления школой плавания.",
            "Клиенты, расписание, посещаемость, абонементы, платежи, отчеты и коммуникации собраны вокруг единой логики.",
            "Главная ценность — меньше ошибок, быстрее контроль оплат, прозрачная история и готовая база для масштабирования.",
        ],
        "image": "summary.png",
    },
]


def font(size=34, bold=False):
    candidates = [
        r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
        r"C:\Windows\Fonts\arialbd.ttf" if bold else r"C:\Windows\Fonts\arial.ttf",
    ]
    for candidate in candidates:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size=size)
    return ImageFont.load_default()


def draw_round_rect(draw, box, radius, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def text_wrap(draw, text, fnt, max_width):
    words = text.split()
    lines = []
    current = ""
    for word in words:
        trial = word if not current else f"{current} {word}"
        if draw.textbbox((0, 0), trial, font=fnt)[2] <= max_width:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)
    return lines


def draw_text_box(draw, xy, text, fnt, fill, max_width, line_gap=8):
    x, y = xy
    for line in text_wrap(draw, text, fnt, max_width):
        draw.text((x, y), line, font=fnt, fill=fill)
        y += fnt.size + line_gap
    return y


def new_canvas():
    img = Image.new("RGB", (1600, 900), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 1600, 900), fill=(255, 255, 255))
    draw.rectangle((0, 0, 1600, 12), fill=(0, 129, 146))
    return img, draw


def save_image(name, painter):
    ASSETS.mkdir(parents=True, exist_ok=True)
    img, draw = new_canvas()
    painter(draw)
    path = ASSETS / name
    img.save(path)
    return path


def create_visuals():
    title = font(52, True)
    h2 = font(34, True)
    h3 = font(28, True)
    body = font(24)
    small = font(20)
    tiny = font(16)

    def cover(draw):
        draw_round_rect(draw, (870, 120, 1470, 720), 28, (245, 250, 252), (210, 225, 232), 2)
        draw_round_rect(draw, (920, 175, 1420, 270), 18, (226, 247, 250), None)
        draw.text((950, 200), "Сегодня", font=h3, fill=(25, 38, 52))
        for i, (label, value, color) in enumerate([
            ("Активные клиенты", "486", (0, 129, 146)),
            ("Занятий на неделе", "78", (54, 113, 222)),
            ("К оплате", "12", (221, 117, 0)),
        ]):
            x = 925 + i * 160
            draw_round_rect(draw, (x, 320, x + 135, 450), 18, (255, 255, 255), (216, 226, 232), 2)
            draw.text((x + 20, 342), value, font=h2, fill=color)
            draw_text_box(draw, (x + 20, 392), label, tiny, (73, 84, 97), 95, 2)
        for y, txt in [(520, "17:00  Группа A  |  Pool 1"), (585, "18:00  Индивидуально  |  Coach Anna"), (650, "19:00  Группа B  |  Pool 2")]:
            draw_round_rect(draw, (925, y, 1415, y + 44), 12, (255, 255, 255), (216, 226, 232), 1)
            draw.text((950, y + 10), txt, font=small, fill=(25, 38, 52))
        draw_round_rect(draw, (100, 170, 690, 610), 32, (226, 247, 250), None)
        draw.arc((155, 245, 355, 445), 200, 340, fill=(0, 129, 146), width=14)
        draw.arc((285, 310, 535, 560), 200, 340, fill=(54, 113, 222), width=14)
        draw.ellipse((185, 365, 230, 410), fill=(0, 129, 146))
        draw.ellipse((445, 465, 490, 510), fill=(54, 113, 222))

    def about(draw):
        roles = [("Админ", "управляет клиентами, оплатами и расписанием"), ("Тренер", "видит занятия и отмечает посещаемость"), ("Родитель", "контролирует абонемент и платежи")]
        for i, (name, desc) in enumerate(roles):
            x = 120 + i * 470
            draw_round_rect(draw, (x, 180, x + 390, 610), 30, (245, 250, 252), (216, 226, 232), 2)
            draw.ellipse((x + 135, 235, x + 255, 355), fill=(226, 247, 250), outline=(0, 129, 146), width=4)
            draw.text((x + 150, 275), name[:2], font=h2, fill=(0, 91, 106))
            draw.text((x + 55, 410), name, font=h2, fill=(25, 38, 52))
            draw_text_box(draw, (x + 55, 465), desc, body, (73, 84, 97), 280)
        draw.text((120, 720), "Одна база вместо разрозненных таблиц, чатов и ручных сверок.", font=h2, fill=(0, 91, 106))

    def capabilities(draw):
        items = [("Клиенты", "семьи, участники, группы"), ("Расписание", "занятия, тренеры, конфликты"), ("Абонементы", "остатки, заморозки, списания"), ("Финансы", "начисления, платежи, долги"), ("Коммуникации", "шаблоны, согласия, логи")]
        for i, (name, desc) in enumerate(items):
            x = 100 + (i % 3) * 480
            y = 150 + (i // 3) * 250
            draw_round_rect(draw, (x, y, x + 390, y + 175), 26, (255, 255, 255), (216, 226, 232), 2)
            draw_round_rect(draw, (x + 30, y + 32, x + 95, y + 97), 18, (226, 247, 250), None)
            draw.text((x + 50, y + 48), str(i + 1), font=h3, fill=(0, 91, 106))
            draw.text((x + 125, y + 35), name, font=h3, fill=(25, 38, 52))
            draw_text_box(draw, (x + 125, y + 82), desc, small, (73, 84, 97), 220)

    def flow(draw):
        steps = ["Клиент", "Группа", "Абонемент", "Занятие", "Отчет"]
        for i, step in enumerate(steps):
            x = 120 + i * 290
            draw_round_rect(draw, (x, 300, x + 200, 420), 28, (226, 247, 250), (0, 129, 146), 2)
            draw.text((x + 45, 338), step, font=h3, fill=(0, 91, 106))
            if i < len(steps) - 1:
                draw.line((x + 210, 360, x + 280, 360), fill=(0, 129, 146), width=5)
                draw.polygon([(x + 280, 360), (x + 260, 348), (x + 260, 372)], fill=(0, 129, 146))
        draw_text_box(draw, (150, 540), "Результат: CRM автоматически обновляет остатки занятий, задолженности, историю посещений и поводы для напоминаний.", body, (73, 84, 97), 1250)

    def cases(draw):
        labels = [("Продажа", "абонемент + платеж"), ("Посещаемость", "статусы и списания"), ("Долги", "контроль и напоминания"), ("Отчеты", "доходы и посещаемость")]
        for i, (name, desc) in enumerate(labels):
            x = 130 + (i % 2) * 680
            y = 160 + (i // 2) * 270
            draw_round_rect(draw, (x, y, x + 560, y + 190), 28, (245, 250, 252), (216, 226, 232), 2)
            draw.ellipse((x + 35, y + 42, x + 125, y + 132), fill=(0, 129, 146))
            draw.text((x + 67, y + 65), str(i + 1), font=h3, fill=(255, 255, 255))
            draw.text((x + 160, y + 45), name, font=h2, fill=(25, 38, 52))
            draw.text((x + 160, y + 98), desc, font=body, fill=(73, 84, 97))

    def limits(draw):
        items = ["Онлайн-оплата", "Провайдеры SMS/Telegram", "Production scheduler", "Регулярный restore drill"]
        for i, item in enumerate(items):
            y = 170 + i * 140
            draw_round_rect(draw, (170, y, 1330, y + 85), 20, (255, 245, 219), (235, 204, 132), 2)
            draw.text((215, y + 24), "!", font=h2, fill=(221, 117, 0))
            draw.text((285, y + 26), item, font=h3, fill=(25, 38, 52))
        draw_text_box(draw, (170, 760), "Ограничения понятны и закрываются поэтапно: сначала production-процессы, затем каналы коммуникации и коммерческие интеграции.", small, (73, 84, 97), 1150)

    def improvements(draw):
        rows = [("UX", 74), ("Функционал", 69), ("Автоматизация", 78), ("Коммерция", 56)]
        for i, (name, value) in enumerate(rows):
            y = 190 + i * 130
            draw.text((160, y), name, font=h3, fill=(25, 38, 52))
            draw_round_rect(draw, (450, y + 12, 1250, y + 52), 18, (232, 241, 255), None)
            draw_round_rect(draw, (450, y + 12, 450 + value * 8, y + 52), 18, (0, 129, 146), None)
            draw.text((1280, y + 7), f"{value}%", font=body, fill=(0, 91, 106))

    def roadmap(draw):
        x_positions = [140, 590, 1040]
        titles = ["Quick wins", "Среднесрочно", "Продвинутые"]
        fills = [(226, 247, 250), (232, 241, 255), (231, 248, 239)]
        for i, (x, name) in enumerate(zip(x_positions, titles)):
            draw_round_rect(draw, (x, 180, x + 340, 650), 30, fills[i], (216, 226, 232), 2)
            draw.text((x + 40, 235), name, font=h3, fill=(25, 38, 52))
            draw.line((x + 40, 300, x + 300, 300), fill=(0, 129, 146), width=4)
            draw_text_box(draw, (x + 40, 340), SLIDES[7]["cards"][i][1], small, (73, 84, 97), 260, 6)

    def value(draw):
        bars = [("Время", 78, (0, 129, 146)), ("Деньги", 66, (54, 113, 222)), ("Контроль", 86, (18, 148, 88))]
        for i, (name, val, color) in enumerate(bars):
            x = 220 + i * 420
            draw_round_rect(draw, (x, 220, x + 260, 640), 28, (245, 250, 252), (216, 226, 232), 2)
            draw.rectangle((x + 75, 575 - val * 4, x + 185, 575), fill=color)
            draw.text((x + 78, 600), name, font=h3, fill=(25, 38, 52))
            draw.text((x + 90, 175), f"{val}%", font=h2, fill=color)
        draw.text((170, 735), "Ценность CRM — в управляемости: меньше ручной работы, быстрее реакция, прозрачнее финансы.", font=body, fill=(73, 84, 97))

    def summary(draw):
        draw_round_rect(draw, (185, 180, 1415, 650), 34, (245, 250, 252), (216, 226, 232), 2)
        draw.text((260, 250), "Готовая база для pilot launch", font=title, fill=(0, 91, 106))
        draw_text_box(draw, (265, 345), "SwimCRM закрывает ежедневную операционную работу школы и дает понятный путь развития: production-процессы, коммуникации, оплаты и аналитика.", body, (73, 84, 97), 960)
        draw_round_rect(draw, (265, 535, 625, 600), 22, (0, 129, 146), None)
        draw.text((300, 550), "Следующий шаг: pilot", font=h3, fill=(255, 255, 255))

    painters = {
        "cover.png": cover,
        "about.png": about,
        "capabilities.png": capabilities,
        "flow.png": flow,
        "cases.png": cases,
        "limits.png": limits,
        "improvements.png": improvements,
        "roadmap.png": roadmap,
        "value.png": value,
        "summary.png": summary,
    }
    return {name: save_image(name, painter) for name, painter in painters.items()}


def set_slide_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_textbox(slide, left, top, width, height, text, size=18, color=BODY, bold=False, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    if align:
        paragraph.alignment = align
    for run in paragraph.runs:
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_header(slide, number, title):
    add_textbox(slide, Inches(0.55), Inches(0.30), Inches(0.6), Inches(0.35), f"{number:02}", 12, TEAL, True)
    add_textbox(slide, Inches(1.15), Inches(0.22), Inches(9.2), Inches(0.55), title, 28, NAVY, True)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.90), Inches(11.95), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def add_card(slide, left, top, width, height, title, body, fill=TEAL_SOFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    add_textbox(slide, left + Inches(0.24), top + Inches(0.16), width - Inches(0.48), Inches(0.32), title, 13, NAVY, True)
    add_textbox(slide, left + Inches(0.24), top + Inches(0.54), width - Inches(0.48), height - Inches(0.66), body, 9, BODY)


def add_bullets(slide, left, top, width, items, size=15):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.11), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        add_textbox(slide, left + Inches(0.25), y, width - Inches(0.25), Inches(0.72), item, size, BODY)
        y += Inches(0.88)


def build_deck(visuals):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Cover
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide)
    slide.shapes.add_picture(str(visuals["cover.png"]), Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    add_textbox(slide, Inches(0.65), Inches(0.55), Inches(6.5), Inches(0.5), SLIDES[0]["tag"], 13, TEAL_DARK, True)
    add_textbox(slide, Inches(0.65), Inches(1.25), Inches(6.7), Inches(1.15), SLIDES[0]["title"], 42, NAVY, True)
    add_textbox(slide, Inches(0.70), Inches(2.55), Inches(5.5), Inches(1.1), SLIDES[0]["subtitle"], 18, BODY)

    for idx, data in enumerate(SLIDES[1:], start=2):
        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide)
        add_header(slide, idx, data["title"])

        if data.get("image"):
            slide.shapes.add_picture(str(visuals[data["image"]]), Inches(6.85), Inches(1.18), Inches(5.55), Inches(3.12))

        if "body" in data:
            add_bullets(slide, Inches(0.80), Inches(1.25), Inches(5.55), data["body"], 12)
        if "cards" in data:
            fills = [TEAL_SOFT, BLUE_SOFT, GREEN_SOFT, AMBER_SOFT]
            for card_idx, (card_title, card_body) in enumerate(data["cards"]):
                left = Inches(0.75)
                top = Inches(1.18 + card_idx * 1.25)
                add_card(slide, left, top, Inches(5.55), Inches(1.03), card_title, card_body, fills[card_idx % len(fills)])

        footer = add_textbox(slide, Inches(0.70), Inches(6.93), Inches(3.2), Inches(0.25), "SwimCRM | client presentation", 8, MUTED)
        footer.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


def main():
    visuals = create_visuals()
    build_deck(visuals)
    print(OUT)


if __name__ == "__main__":
    main()
